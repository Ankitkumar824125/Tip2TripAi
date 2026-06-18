import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Installing python-pptx...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Define color scheme
    COLOR_PRIMARY = RGBColor(23, 128, 61)   # Brand Green
    COLOR_SECONDARY = RGBColor(220, 38, 38) # Brand Red
    COLOR_DARK = RGBColor(15, 23, 42)       # Text Slate
    COLOR_GRAY = RGBColor(100, 116, 139)     # Text Gray
    COLOR_LIGHT = RGBColor(247, 245, 240)   # Light Beige

    # ------------------ SLIDE 1: TITLE SLIDE ------------------
    slide_layout = prs.slide_layouts[5] # Blank layout with title placeholder
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT

    # Title Box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(8.0), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Tip2Trip AI Travel Planner"
    p.font.name = "Arial"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Four-Week Internship & Progress Report"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_DARK
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Comprehensive Summary: Weeks 01 to 04"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_SECONDARY
    p3.alignment = PP_ALIGN.CENTER

    p4 = tf.add_paragraph()
    p4.text = "\nPrepared for: Internship Evaluation / Faculty Review"
    p4.font.name = "Arial"
    p4.font.size = Pt(12)
    p4.font.italic = True
    p4.font.color.rgb = COLOR_GRAY
    p4.alignment = PP_ALIGN.CENTER

    # Helper function to add structured slide
    def add_content_slide(title_text, points_list):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
        
        # Background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_LIGHT
        
        # Title Box
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(8.5), Inches(0.8))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY
        
        # Content Box
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(8.5), Inches(5.5))
        tf_content = content_box.text_frame
        tf_content.word_wrap = True
        
        for idx, item in enumerate(points_list):
            p_bullet = tf_content.paragraphs[0] if idx == 0 else tf_content.add_paragraph()
            p_bullet.space_after = Pt(10)
            
            # Highlight keyword
            run_title = p_bullet.add_run()
            run_title.text = item[0] + ": "
            run_title.font.name = "Arial"
            run_title.font.size = Pt(16)
            run_title.font.bold = True
            run_title.font.color.rgb = COLOR_DARK
            
            # Sub text
            run_desc = p_bullet.add_run()
            run_desc.text = item[1]
            run_desc.font.name = "Arial"
            run_desc.font.size = Pt(15)
            run_desc.font.color.rgb = COLOR_GRAY
            
            # Code reference if available
            if len(item) > 2:
                run_ref = p_bullet.add_run()
                run_ref.text = f"\n  (File: {item[2]})"
                run_ref.font.name = "Courier New"
                run_ref.font.size = Pt(12)
                run_ref.font.italic = True
                run_ref.font.color.rgb = COLOR_PRIMARY

    # ------------------ SLIDE 2: WEEK 1 - PROJECT GOALS ------------------
    add_content_slide(
        "Week 01: Project Goals & Objectives",
        [
            ("Project Vision", "Understood the core platform objective to provide seamless AI itineraries, travel buddies finder, and expense splitting.", "App.jsx"),
            ("Workflow Analysis", "Analyzed the end-to-end traveler workflow: Hero Search -> Feature Grid -> Buddy Match -> Timeline Dashboard.", "src/components/"),
            ("Routing Architecture", "Mapped the component navigation structure, tracking page visibility across conditional routes.", "Navbar.jsx & App.jsx")
        ]
    )

    # ------------------ SLIDE 3: WEEK 1 - TECHNICAL STRUCTURE ------------------
    add_content_slide(
        "Week 01: Technical Structure & Dev Setup",
        [
            ("Frontend & Backend split", "Modular React component setup located in src/components/ coupled with Express Node backend index.js.", "src/ & server/"),
            ("Data Store", "Local JSON-based mockup database used for data persistence of trips, users, and community posts.", "server/db.json"),
            ("Environment Setup", "Configured and ran Vite dev server on http://localhost:5173 and Express server on http://localhost:5000.", "vite.config.js & package.json")
        ]
    )

    # ------------------ SLIDE 4: WEEK 2 - ROUTING & NAVIGATION ------------------
    add_content_slide(
        "Week 02: Routing & Component Flow",
        [
            ("State-based Routing", "Used react state variables to transition between views (home, profile, signup, login) smoothly.", "App.jsx"),
            ("Route Protection", "Implemented auth-guards that automatically intercept unauthenticated requests and redirect users to login.", "App.jsx"),
            ("Dynamic Transitions", "Managed transitions and visual feedback overlays during screen updates.", "App.jsx")
        ]
    )

    # ------------------ SLIDE 5: WEEK 2 - NAVBAR MODIFICATIONS ------------------
    add_content_slide(
        "Week 02: Navbar Customization & Brand Assets",
        [
            ("Conditional Menu Items", "Modified menus to display customized items (Avatar, Profile, Sign Out) only when user is logged in.", "src/components/Navbar.jsx"),
            ("Navigation Callback", "Bound the onNavigate event listener to trigger currentPage updates in the parent component.", "src/components/Navbar.jsx"),
            ("Branding Icons", "Updated brand assets, including the compass logo icon and custom text alignments.", "src/components/Navbar.jsx")
        ]
    )

    # ------------------ SLIDE 6: WEEK 3 - LANDING PAGE MODULES ------------------
    add_content_slide(
        "Week 03: Landing Page Modules & Layout updates",
        [
            ("Hero Search Widget", "Updated search input fields, incorporating destination name fields, date selectors, and budget type filters.", "src/components/Hero.jsx"),
            ("Features Grid Layout", "Refactored the capabilities list into grid-based styled cards presenting the main visual benefits.", "src/components/Features.jsx"),
            ("Stepper Roadmap", "Improved visual timeline walkthrough indicators (Select, Personalize, Connect) to guide first-time travelers.", "src/components/HowItWorks.jsx"),
            ("Responsive Alignments", "Refactored grid mappings and flex arrangements to ensure complete alignment on mobile, tablet, and widescreen layouts.", "App.css")
        ]
    )

    # ------------------ SLIDE 7: WEEK 4 - TRAVEL BUDDY MATCHER ------------------
    add_content_slide(
        "Week 04: Travel Buddy Vibe-Checker",
        [
            ("User Profiles & Cards", "Rendered traveler card components showing matching profiles, compatible match percentages, avatar status glows, and bio texts.", "src/components/TravelBuddy.jsx"),
            ("Destination Filters", "Configured action buttons to filter matches dynamically by targeted locations or query strings.", "src/components/TravelBuddy.jsx"),
            ("Like & Connect Handles", "Integrated Express API callbacks (/api/buddies) to support persistent client state for liking or connecting with buddies.", "src/components/TravelBuddy.jsx")
        ]
    )

    # ------------------ SLIDE 8: WEEK 4 - COMMUNITY BROADCAST FEED ------------------
    add_content_slide(
        "Week 04: Community Broadcast Feed",
        [
            ("Post Creator Card", "Created a broadcast card that lets users write updates, add a location tag, and submit them in real-time.", "src/components/CommunityFeed.jsx"),
            ("Social Timelines", "Rendered interactive timeline posts, including user information tags, location identifiers, and attached custom media.", "src/components/CommunityFeed.jsx"),
            ("Social Likes & Feedback", "Wired REST API endpoints (/api/posts) to fetch updates, trigger like updates, and report comment metrics.", "src/components/CommunityFeed.jsx")
        ]
    )

    # ------------------ SLIDE 9: KEY TAKEAWAYS ------------------
    add_content_slide(
        "Tip2Trip: Key Takeaways & Intern Progress",
        [
            ("Full Stack Connection", "Successfully created front-to-back connection, fetching assets asynchronously via REST APIs.", "src/ & server/"),
            ("Elegant Handling", "Implemented custom error catches in signup/login pages to protect against server/network issues.", "Login.jsx & Signup.jsx"),
            ("Visual Brand Value", "Established a premium dark/light mode layout using tailored glassmorphism design styles.", "App.css")
        ]
    )

    # Save presentation
    output_path = "Tip2Trip_Full_Report.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
